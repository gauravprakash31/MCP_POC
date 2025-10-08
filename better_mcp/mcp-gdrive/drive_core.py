# mcp-gdrive/drive_core.py
import io, pathlib
from typing import Dict, Any, Tuple, List
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials

import chardet
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
import pandas as pd

ROOT = pathlib.Path(__file__).resolve().parent
TOKEN_PATH = ROOT / "token.json"
SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]

GOOGLE_APPS_PREFIX = "application/vnd.google-apps"
EXPORT_MAP = {
    "application/vnd.google-apps.document": "text/plain",
    "application/vnd.google-apps.spreadsheet": "text/csv",
    "application/vnd.google-apps.presentation": "text/plain",
}

def _creds() -> Credentials:
    if not TOKEN_PATH.exists():
        raise RuntimeError("Missing token.json. Run auth_setup.py.")
    c = Credentials.from_authorized_user_file(str(TOKEN_PATH), SCOPES)
    if not c.valid:
        if c.expired and c.refresh_token:
            c.refresh(Request()); TOKEN_PATH.write_text(c.to_json())
        else:
            raise RuntimeError("Re-run auth_setup.py.")
    return c

def _drive():
    return build("drive", "v3", credentials=_creds(), cache_discovery=False)

def search_files(query: str, page_size: int = 25) -> List[Dict[str, Any]]:
    svc = _drive()
    resp = svc.files().list(
        q=query, pageSize=page_size,
        fields="files(id,name,mimeType,modifiedTime,webViewLink)"
    ).execute()
    return resp.get("files", [])

def get_metadata(file_id: str) -> Dict[str, Any]:
    svc = _drive()
    return svc.files().get(
        fileId=file_id,
        fields="id,name,mimeType,size,owners,modifiedTime,parents,webViewLink"
    ).execute()

def _export_google_file(file_id: str, export_mime: str) -> bytes:
    svc = _drive()
    return svc.files().export(fileId=file_id, mimeType=export_mime).execute()

def _download_binary(file_id: str) -> Tuple[bytes, str]:
    svc = _drive()
    meta = svc.files().get(fileId=file_id, fields="mimeType").execute()
    mime = meta["mimeType"]
    req = svc.files().get_media(fileId=file_id)
    buf = io.BytesIO()
    downloader = MediaIoBaseDownload(buf, req)
    done = False
    while not done:
        status, done = downloader.next_chunk()
    return buf.getvalue(), mime

def get_file_bytes(file_id: str) -> Tuple[bytes, str]:
    meta = get_metadata(file_id)
    mime = meta["mimeType"]
    if mime.startswith(GOOGLE_APPS_PREFIX):
        export_mime = EXPORT_MAP.get(mime, "application/pdf")
        return _export_google_file(file_id, export_mime), export_mime
    return _download_binary(file_id)

# --- extraction ---
def _detect_text_bytes(b: bytes) -> str:
    enc = chardet.detect(b).get("encoding") or "utf-8"
    try: return b.decode(enc, errors="replace")
    except Exception: return b.decode("utf-8", errors="replace")

def _pdf_to_text(b: bytes) -> str:
    reader = PdfReader(io.BytesIO(b))
    return "\n".join([(p.extract_text() or "") for p in reader.pages]).strip() or "[No extractable text in PDF]"

def _docx_to_text(b: bytes) -> str:
    with io.BytesIO(b) as bio: doc = DocxDocument(bio)
    return "\n".join([p.text for p in doc.paragraphs]).strip()

def _pptx_to_text(b: bytes) -> str:
    with io.BytesIO(b) as bio: pres = PptxPresentation(bio)
    texts = []
    for slide in pres.slides:
        for sh in slide.shapes:
            if hasattr(sh, "text"): texts.append(sh.text)
    return "\n".join(texts).strip() or "[No text found in slides]"

def _xlsx_to_text(b: bytes) -> str:
    with io.BytesIO(b) as bio: xl = pd.ExcelFile(bio)
    parts = []
    for name in xl.sheet_names:
        df = xl.parse(name).head(50)
        parts.append(f"# Sheet: {name}\n" + df.to_csv(index=False))
    return "\n\n".join(parts)

def bytes_to_text(b: bytes, mime: str) -> str:
    mime = (mime or "").lower()
    if mime == "application/pdf": return _pdf_to_text(b)
    if mime in ("text/plain","text/markdown","application/json","text/html","application/vnd.google-apps.script+json"):
        return _detect_text_bytes(b)
    if mime in ("text/csv","text/tab-separated-values"): return _detect_text_bytes(b)
    if mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document": return _docx_to_text(b)
    if mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation": return _pptx_to_text(b)
    if mime == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": return _xlsx_to_text(b)
    txt = _detect_text_bytes(b)
    return txt if txt.strip() else f"[Unsupported or binary type: {mime or 'unknown'}]"

def get_text(file_id: str, max_chars: int = 60_000) -> str:
    b, mime = get_file_bytes(file_id)
    t = bytes_to_text(b, mime)
    return (t[:max_chars] + "\n[Truncated]") if len(t) > max_chars else t
